"""
Chunking + embedding + vector retrieval, using Chroma (local, free, no API key).
Mirrors the ingestion -> chunking -> embeddings -> vector search pattern used
in the Clinical Decision Support project, but fully open-source/local here.
"""
import hashlib
import os
import re
import base64
import mimetypes
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

import chromadb

CHROMA_DIR = os.path.join(os.path.dirname(__file__), "..", ".chroma")
VECTOR_DIM = 256
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".mp4", ".webm"}


class HashingEmbeddingFunction(chromadb.EmbeddingFunction):
    """
    Lightweight, dependency-free embedding function: hashes each word into
    one of VECTOR_DIM buckets (the "hashing trick"), giving a bag-of-words
    vector with no model download and no torch/onnx dependency. Good
    enough to demo retrieval end-to-end offline.

    Swap this for OpenAI/Cohere embeddings or sentence-transformers in a
    real deployment for much better semantic retrieval quality — the
    reason to call that out explicitly in an interview is that it shows
    you made a deliberate tradeoff (portability/zero-cost demo vs.
    retrieval quality), not that you didn't know better options exist.
    """

    def name(self) -> str:
        return "hashing-bow-v1"

    def __call__(self, input: list[str]) -> list[list[float]]:
        return [self._embed(text) for text in input]

    def _embed(self, text: str) -> list[float]:
        vec = [0.0] * VECTOR_DIM
        words = re.findall(r"[a-z0-9]+", text.lower())
        for w in words:
            idx = int(hashlib.md5(w.encode()).hexdigest(), 16) % VECTOR_DIM
            vec[idx] += 1.0
        norm = sum(v * v for v in vec) ** 0.5
        return [v / norm for v in vec] if norm > 0 else vec


def _get_collection():
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    return client.get_or_create_collection(name="docs", embedding_function=HashingEmbeddingFunction())


def chunk_text(text: str, chunk_size: int = 400, overlap: int = 50) -> list[str]:
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start += chunk_size - overlap
    return chunks


def _extract_docx_text(file_path: str) -> str:
    from docx import Document

    document = Document(file_path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(paragraphs)


def _extract_pptx_text(file_path: str) -> str:
    from pptx import Presentation

    presentation = Presentation(file_path)
    slides = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if text:
            slides.append(f"Slide {slide_number}: " + "\n".join(text))
    return "\n".join(slides)


def _get_groq_client():
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        raise ValueError("GROQ_API_KEY is required to process audio and images.")

    from groq import Groq

    return Groq(api_key=api_key)


def _transcribe_audio(file_path: str) -> str:
    client = _get_groq_client()
    with open(file_path, "rb") as audio_file:
        transcription = client.audio.transcriptions.create(
            file=audio_file,
            model=os.getenv("GROQ_AUDIO_MODEL", "whisper-large-v3-turbo"),
            response_format="text",
        )
    return str(transcription)


def _describe_image(file_path: str) -> str:
    client = _get_groq_client()
    mime_type = mimetypes.guess_type(file_path)[0] or "image/jpeg"
    with open(file_path, "rb") as image_file:
        image_data = base64.b64encode(image_file.read()).decode("ascii")

    response = client.chat.completions.create(
        model=os.getenv("GROQ_VISION_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct"),
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Describe this image accurately. Extract visible text, charts, labels, and key facts for a retrieval system."},
                    {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}},
                ],
            }
        ],
        temperature=0,
        max_tokens=800,
    )
    return response.choices[0].message.content or ""


def extract_text(file_path: str) -> str:
    """Extract searchable text from a supported document, audio, or image file."""
    extension = os.path.splitext(file_path)[1].lower()
    if extension == ".txt":
        with open(file_path, "r", encoding="utf-8") as text_file:
            return text_file.read()
    if extension == ".pdf" and HAS_PDF:
        with open(file_path, "rb") as pdf_file:
            return "\n".join(page.extract_text() or "" for page in PyPDF2.PdfReader(pdf_file).pages)
    if extension == ".docx":
        return _extract_docx_text(file_path)
    if extension == ".pptx":
        return _extract_pptx_text(file_path)
    if extension in AUDIO_EXTENSIONS:
        return _transcribe_audio(file_path)
    if extension in IMAGE_EXTENSIONS:
        return _describe_image(file_path)
    return ""


def ingest_file(file_path: str) -> int:
    """Ingest a supported document, audio recording, or image file."""
    collection = _get_collection()
    count = 0
    
    if not os.path.isfile(file_path):
        return 0
    
    fname = os.path.basename(file_path)
    try:
        text = extract_text(file_path)
    except Exception:
        return 0
    
    if text and text.strip():
        for i, chunk in enumerate(chunk_text(text)):
            collection.add(
                documents=[chunk],
                ids=[f"{fname}-{i}"],
                metadatas=[{"source": fname, "chunk": i, "modality": os.path.splitext(fname)[1].lower()}],
            )
            count += 1
    
    return count


def ingest_directory(directory: str) -> int:
    """Ingest all supported files within a directory."""
    count = 0
    
    if not os.path.isdir(directory):
        return 0
    
    for fname in os.listdir(directory):
        fpath = os.path.join(directory, fname)
        if not os.path.isfile(fpath):
            continue
            
        count += ingest_file(fpath)
    
    return count


def retrieve(query: str, top_k: int = 3) -> list[dict]:
    collection = _get_collection()
    if collection.count() == 0:
        return []
    results = collection.query(query_texts=[query], n_results=min(top_k, collection.count()))
    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    return [{"text": d, "source": m.get("source")} for d, m in zip(docs, metas)]
